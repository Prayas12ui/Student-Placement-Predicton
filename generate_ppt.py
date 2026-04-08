from pptx import Presentation
from pptx.util import Inches, Pt
import os

def create_presentation():
    # Create presentation
    prs = Presentation()

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Student Performance & Placement Analytics"
    subtitle.text = "Predicting Placement Outcomes Using Machine Learning"

    # Slide 1: Introduction
    bullet_slide_layout = prs.slide_layouts[1]
    slide1 = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide1.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Introduction & Objective"
    tf = body_shape.text_frame
    tf.text = "Objective: Analyze student performance to predict placement outcomes."
    p = tf.add_paragraph()
    p.text = "Goal is to identify key factors influencing placements (CGPA, Skills, Projects)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Provides actionable insights for students and institutions."
    p.level = 1

    # Slide 2: Dataset Overview
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide2.shapes.title
    body_shape = slide2.placeholders[1]
    title_shape.text = "Dataset Features"
    tf = body_shape.text_frame
    tf.text = "The dataset consists of 1000 records containing:"
    features = [
        "10th and 12th Marks (%)",
        "B.Tech CGPA",
        "Programming & Communication Skills (Score)",
        "Number of Internships & Projects",
        "Active Backlogs",
        "Placement Status (Target)"
    ]
    for feat in features:
        p = tf.add_paragraph()
        p.text = feat
        p.level = 1

    # Slide 3: Visualization 1
    blank_slide_layout = prs.slide_layouts[5] # Title Only
    slide3 = prs.slides.add_slide(blank_slide_layout)
    slide3.shapes.title.text = "Data Analysis: CGPA vs Placement"
    img_path1 = "plots/cgpa_vs_placement.png"
    if os.path.exists(img_path1):
        slide3.shapes.add_picture(img_path1, Inches(1), Inches(2), width=Inches(8))
    else:
        txBox = slide3.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        txBox.text_frame.text = "[Image cgpa_vs_placement.png not found. Run train_model.py first]"

    # Slide 4: Visualization 2
    slide4 = prs.slides.add_slide(blank_slide_layout)
    slide4.shapes.title.text = "Data Analysis: Skills vs Placement"
    img_path2 = "plots/skills_vs_placement.png"
    if os.path.exists(img_path2):
        slide4.shapes.add_picture(img_path2, Inches(1), Inches(2), width=Inches(8))
    else:
        txBox = slide4.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        txBox.text_frame.text = "[Image skills_vs_placement.png not found. Run train_model.py first]"

    # Slide 5: Machine Learning Model
    slide5 = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide5.shapes.title
    body_shape = slide5.placeholders[1]
    title_shape.text = "Machine Learning Implementation"
    tf = body_shape.text_frame
    tf.text = "Algorithm: Random Forest Classifier"
    p = tf.add_paragraph()
    p.text = "Ensemble method used for robust predictions and avoiding overfitting."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Training / Test Split: 80% / 20%"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Model achieves high accuracy in classifying 'Placed' and 'Not Placed' students."
    p.level = 1

    # Slide 6: Conclusion
    slide6 = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide6.shapes.title
    body_shape = slide6.placeholders[1]
    title_shape.text = "Conclusion & Future Work"
    tf = body_shape.text_frame
    tf.text = "Key Findings:"
    p = tf.add_paragraph()
    p.text = "CGPA and Programming Skills are strong predictors of placement."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Projects and Internships significantly boost placement chances."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Future Work:"
    p = tf.add_paragraph()
    p.text = "Integration with live college databases."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Deploying the Streamlit dashboard to a public cloud."
    p.level = 1

    prs.save('Placement_Project_Presentation.pptx')
    print("Presentation saved as Placement_Project_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
